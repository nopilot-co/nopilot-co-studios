#!/usr/bin/env python3
"""Editable PPTX engine (#19). Standalone; run:
design/.venv/bin/python tests/test_pptx.py
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import pptx_render  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


TOK = {
    "color": {
        "neutral": "#142F54",
        "surface": "#21456b",
        "tertiary": "#B14B33",
        "on_primary": "#FFFFFF",
        "secondary": "#6B7280",
        "primary": "#2A3548",
    },
    "space": {"sm": "8pt", "md": "16pt", "lg": "32pt"},
    "radius": {"sm": "2pt", "md": "4pt", "lg": "8pt"},
}

# split_slides: headings and cover/section blocks start new slides.
md = (
    "::: cover-slide\n# Q3 Strategy\nSubtitle here.\n:::\n\n"
    "## Context\nA point about the market.\n\n"
    "## Plan\n- First\n- Second\n"
)
slides = pptx_render.split_slides(md)
check("split: 3 slides", len(slides) == 3, str(len(slides)))
check("split: first is cover", slides[0]["kind"] == "cover-slide", str(slides[0]))

# build_pptx writes a real .pptx with the right slide count.
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "deck.pptx"
    pptx_render.build_pptx(md, TOK, out)
    check("pptx written", out.exists() and out.stat().st_size > 0)
    from pptx import Presentation

    prs = Presentation(str(out))
    check("pptx slide count", len(list(prs.slides)) == 3, str(len(list(prs.slides))))

# Tier 2: a slide bearing ::: kpi / ::: panel / a markdown table emits native shapes.
md2 = (
    "## Numbers\n\n"
    "::: kpi\n87% faster\n:::\n\n"
    "::: panel\nA framed aside.\n:::\n\n"
    "| Metric | Value |\n|---|---|\n| Revenue | 1.8M |\n| Margin | 17% |\n"
)
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d2.pptx"
    pptx_render.build_pptx(md2, TOK, out)
    from pptx import Presentation

    prs = Presentation(str(out))
    sl = list(prs.slides)[0]
    has_table = any(getattr(sh, "has_table", False) for sh in sl.shapes)
    has_auto = any("AUTO_SHAPE" in str(sh.shape_type) for sh in sl.shapes)
    check(
        "tier2 native table", has_table, str([str(sh.shape_type) for sh in sl.shapes])
    )
    check("tier2 panel autoshape", has_auto)

# Tier 3: ::: chart -> a NATIVE editable chart object.
md3 = "## Revenue\n\n::: chart\ntype: bar\nx: [Q1, Q2, Q3]\ny: [12, 18, 15]\n:::\n"
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d3.pptx"
    pptx_render.build_pptx(md3, TOK, out)
    from pptx import Presentation

    prs = Presentation(str(out))
    sl = list(prs.slides)[0]
    check(
        "tier3 native chart",
        any(getattr(sh, "has_chart", False) for sh in sl.shapes),
        str([str(sh.shape_type) for sh in sl.shapes]),
    )
# bad chart YAML -> slide still builds (degrade)
md3b = "## Oops\n\n::: chart\ntype: nope\n:::\n"
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d3b.pptx"
    pptx_render.build_pptx(md3b, TOK, out)  # must not raise
    check("tier3 bad chart no crash", out.exists())

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print("PASS: pptx (tier 1)")
