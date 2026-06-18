#!/usr/bin/env python3
"""Cross-backend flow parity (ADR-006 / #129). Standalone; run: python3 tests/test_archetype_flow.py

ONE FlowNode reconciles the flow dialects — gslide [{title,caption}], html
"Title — caption" strings, pptx/diagrams nodes:[labels] — and gslide / pptx /
UDS-HTML each render every stage natively (wrapping, never truncating).
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import archetype_ir as ir  # noqa: E402
from studio import gslide, pptx_render, uds_html  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


# 1. the dialects reconcile to the same ordered steps
gsl = ir.normalise_flow([{"title": "Discover", "caption": "Audit"}, {"title": "Design", "caption": "Wireframes"}, {"title": "Build", "caption": "Ship"}])
htmd = ir.normalise_flow({"steps": ["Discover — Audit", "Design — Wireframes", "Build — Ship"]})
ppd = ir.normalise_flow({"nodes": ["Discover", "Design", "Build"]})
check("gslide dialect → 3 steps", [s.title for s in gsl.steps] == ["Discover", "Design", "Build"], str(gsl.steps))
check("html string dialect → title+caption split", [(s.title, s.caption) for s in htmd.steps] == [("Discover", "Audit"), ("Design", "Wireframes"), ("Build", "Ship")], str(htmd.steps))
check("pptx nodes dialect → titles", [s.title for s in ppd.steps] == ["Discover", "Design", "Build"], str(ppd.steps))
check("malformed → FlowNode, no crash", isinstance(ir.normalise_flow("::: not: [yaml"), ir.FlowNode))
check("capabilities[flow] = 3 backends", ir.CAPABILITIES["flow"] == {"gslide", "pptx", "html"})
check("process aliases to flow (gslide/pptx sense)", ir.canonical("process") == "flow")

FLOW_MD = ("---\ntitle: T\n---\n\n## Process {#p}\n\n::: flow\n"
           "- {title: Discover, caption: Audit and interviews}\n"
           "- {title: Design, caption: Wireframes and prototypes}\n"
           "- {title: Build, caption: Implementation and QA}\n:::\n")
TOK = {"color": {"neutral": "#142F54", "surface": "#21456b", "tertiary": "#B14B33",
                 "on_primary": "#FFFFFF", "on_surface": "#FFFFFF", "secondary": "#6B7280", "primary": "#2A3548"}}

# 2. gslide: native numbered chips, every stage
with tempfile.TemporaryDirectory() as td:
    src = Path(td) / "d.md"
    src.write_text(FLOW_MD)
    _t, reqs = gslide.build_requests(src, brand="nopilot")
    inserts = [r["insertText"]["text"] for r in reqs if "insertText" in r]
    check("gslide: stage titles present", all(t in inserts for t in ("Discover", "Design", "Build")), str(inserts))
    check("gslide: numbered 1..3", all(str(i) in inserts for i in (1, 2, 3)))
    check("gslide: captions present", any("interviews" in t for t in inserts))

# 3. pptx: native boxes — now from the {title} dict dialect it could not read before
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d.pptx"
    pptx_render.build_pptx("## Process\n\n::: flow\nsteps:\n  - {title: Discover}\n  - {title: Design}\n  - {title: Build}\n:::\n", TOK, out)
    from pptx import Presentation

    sl = list(Presentation(str(out)).slides)[0]
    autoshapes = [sh for sh in sl.shapes if "AUTO_SHAPE" in str(sh.shape_type)]
    check("pptx: >= 3 flow boxes", len(autoshapes) >= 3, str(len(autoshapes)))

# 4. UDS-HTML: native chips, every stage, wraps
_m, html = uds_html.render_body(FLOW_MD, brand="nopilot")
check("html: uds-flow list", 'class="uds-flow"' in html)
check("html: 3 native steps", html.count("uds-flow__step") == 3, str(html.count("uds-flow__step")))
check("html: titles present", all(t in html for t in ("Discover", "Design", "Build")))
check("html: captions present", "interviews" in html and "Wireframes" in html)
check("html: wraps (flex-wrap)", "flex-wrap:wrap" in html)
check("html: fence consumed", ":::" not in html)

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print("PASS: flow archetype — native parity across gslide / pptx / UDS-HTML")
