#!/usr/bin/env python3
"""Cross-backend chart parity (ADR-006 / #129 — render convergence V0).
Standalone; run: python3 tests/test_archetype_chart.py

Proves: ONE ChartNode normaliser reconciles the canonical (x/y) and gslide
per-bar (series:[{label,value}]) dialects, and gslide / pptx / UDS-HTML each
render that node natively — every bar shown, coloured from the brand ramp.
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import archetype_ir as ir  # noqa: E402
from studio import gslide, pptx_render, uds_html  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


# ---- 1. the normaliser reconciles every dialect into one node --------------
canon = ir.normalise_chart({"type": "bar", "x": ["A", "B", "C"], "y": [3, 9, 6]})
perbar = ir.normalise_chart({"series": [
    {"label": "A", "value": 3}, {"label": "B", "value": 9}, {"label": "C", "value": 6}]})
data_alias = ir.normalise_chart({"data": [
    {"label": "A", "value": 3}, {"label": "B", "value": 9}, {"label": "C", "value": 6}]})
check("canonical categories", canon.categories == ["A", "B", "C"], str(canon.categories))
check("canonical values", canon.series[0].values == [3.0, 9.0, 6.0], str(canon.series[0].values))
check("per-bar → same categories", perbar.categories == ["A", "B", "C"], str(perbar.categories))
check("per-bar → same values (dialects reconcile)", perbar.series[0].values == [3.0, 9.0, 6.0], str(perbar.series[0].values))
check("data: alias == series per-bar", data_alias.series[0].values == perbar.series[0].values)
multi = ir.normalise_chart({"type": "line", "x": ["Q1", "Q2"], "series": [
    {"name": "Plan", "y": [10, 14]}, {"name": "Actual", "y": [12, 18]}]})
check("multi-series kept distinct", len(multi.series) == 2 and multi.series[0].name == "Plan", str([s.name for s in multi.series]))
pie = ir.normalise_chart({"type": "pie", "labels": ["A", "B"], "values": [3, 7]})
check("pie values", pie.series[0].values == [3.0, 7.0], str(pie.series))
check("malformed YAML degrades, no crash", ir.normalise_chart("not: [a, b").is_empty)
check("'30d' display → numeric value", ir.normalise_chart({"series": [{"label": "P0", "value": "30d", "display": "30d"}]}).series[0].values == [30.0])

# ---- 2. capability matrix declares all three backends ----------------------
check("capabilities[chart] = 3 backends", ir.CAPABILITIES["chart"] == {"gslide", "pptx", "html"}, str(ir.CAPABILITIES.get("chart")))

CHART_MD = "---\ntitle: Test\n---\n\n## Numbers {#nums}\n\n::: chart\ntype: bar\nx: [A, B, C]\ny: [3, 9, 6]\n:::\n"
RAMP0 = ir.palette_for("nopilot")[0]  # the brand crimson — proves tokens flow to every backend

# ---- 3. gslide: native bars from the node ----------------------------------
with tempfile.TemporaryDirectory() as td:
    src = Path(td) / "deck.md"
    src.write_text(CHART_MD)
    _title, reqs = gslide.build_requests(src, brand="nopilot")
    bars = [r for r in reqs if "_bar" in r.get("updateShapeProperties", {}).get("objectId", "")]
    inserts = [r["insertText"]["text"] for r in reqs if "insertText" in r]
    check("gslide: 3 native bars", len(bars) == 3, str(len(bars)))
    check("gslide: categories present", all(c in inserts for c in ("A", "B", "C")), str(inserts))
    check("gslide: values present", all(v in inserts for v in ("3", "9", "6")))

    def _near(a: dict, b: dict) -> bool:
        return all(abs(a.get(k, 0) - b.get(k, 0)) <= 0.02 for k in ("red", "green", "blue"))

    bar0 = next((r for r in reqs if r.get("updateShapeProperties", {}).get("objectId", "").endswith("_bar0_0")), None)
    got = bar0["updateShapeProperties"]["shapeProperties"]["shapeBackgroundFill"]["solidFill"]["color"]["rgbColor"] if bar0 else {}
    check("gslide: bar0 uses brand ramp[0]", _near(got, gslide._rgb(RAMP0)), f"{got} vs {RAMP0}")

# ---- 4. pptx: native chart object from the node ----------------------------
TOK = {"color": {"neutral": "#142F54", "surface": "#21456b", "tertiary": "#B14B33",
                 "on_primary": "#FFFFFF", "on_surface": "#FFFFFF", "secondary": "#6B7280", "primary": "#2A3548"}}
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d.pptx"
    pptx_render.build_pptx("## Numbers\n\n::: chart\ntype: bar\nx: [A, B, C]\ny: [3, 9, 6]\n:::\n", TOK, out)
    from pptx import Presentation

    sl = list(Presentation(str(out)).slides)[0]
    check("pptx: native chart object", any(getattr(sh, "has_chart", False) for sh in sl.shapes),
          str([str(sh.shape_type) for sh in sl.shapes]))

# ---- 5. UDS-HTML: native DOM bars from the node ----------------------------
_meta, html = uds_html.render_body(CHART_MD, brand="nopilot")
check("html: uds-chart figure", 'class="uds-chart"' in html)
check("html: 3 native bars", html.count("uds-chart__bar") == 3, str(html.count("uds-chart__bar")))
check("html: value labels present", all(f">{v}<" in html for v in ("3", "9", "6")), "value labels")
check("html: brand ramp colour inlined", RAMP0 in html, RAMP0)
check("html: fence consumed (no ::: leak)", ":::" not in html)

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print("PASS: chart archetype — native parity across gslide / pptx / UDS-HTML")
