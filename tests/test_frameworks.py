#!/usr/bin/env python3
"""Framework renderers (Phase 2) — bullseye, matrix, funnel, heatmap, swimlane,
decision-tree → brand-styled SVG, expanded per export.
Standalone; run: design/.venv/bin/python tests/test_frameworks.py
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import frameworks  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


TOK = {
    "color": {
        "neutral": "#142F54", "surface": "#21456b", "tertiary": "#B14B33",
        "on_primary": "#FFFFFF", "secondary": "#6B7280", "primary": "#2A3548",
    },
}

SPECS = {
    "bullseye": {"rings": [{"ring": "core", "items": ["Alpha", "Beta"]},
                           {"ring": "adjacent", "items": ["Gamma"]}]},
    "matrix": {"axes": {"x": "Effort", "y": "Impact"},
               "items": [{"label": "QuickWin", "x": "low", "y": "high", "quadrant": "Do now"}]},
    "funnel": {"stages": [{"stage": "Visits", "value": 1000}, {"stage": "Signups", "value": 200}]},
    "heatmap-rag": {"rag": True, "rows": ["TeamA"], "cols": ["Q1", "Q2"], "cells": [["green", "red"]]},
    "heatmap-num": {"rows": ["RowA"], "cols": ["X", "Y"], "cells": [[1, 9]]},
    "swimlane": {"lanes": [{"lane": "Cust", "nodes": ["Reqst", "Apprv"]},
                           {"lane": "Ops", "nodes": ["Fulfl"]}]},
    "decision-tree": {"root": "Qstn", "children": [{"condition": "yesp", "root": "Propl"},
                                                   {"condition": "nope", "root": "Nurtr"}]},
}

# each type renders valid SVG containing its key labels.
for key, spec in SPECS.items():
    name = "heatmap" if key.startswith("heatmap") else key
    svg = frameworks.render_svg(name, spec, TOK)
    check(f"{key}: is svg", "<svg" in svg[:600], svg[:80])

# label presence (matplotlib emits text as <text> in this env, per test_charts).
check("bullseye label", "Alpha" in frameworks.render_svg(name="bullseye", spec=SPECS["bullseye"], tokens=TOK))
check("matrix label", "QuickWin" in frameworks.render_svg("matrix", SPECS["matrix"], TOK))
check("funnel stage", "Visits" in frameworks.render_svg("funnel", SPECS["funnel"], TOK))
check("swimlane lane", "Cust" in frameworks.render_svg("swimlane", SPECS["swimlane"], TOK))
dt = frameworks.render_svg("decision-tree", SPECS["decision-tree"], TOK)
check("dtree node + condition", "Propl" in dt and "yesp" in dt, dt[:80])
hm = frameworks.render_svg("heatmap", SPECS["heatmap-rag"], TOK)
check("heatmap rag uses traffic colour", "#2E7D32" in hm or "#2e7d32" in hm.lower(), "green missing")

# bad spec -> render_svg raises (so expand can catch it).
raised = False
try:
    frameworks.render_svg("bullseye", {"rings": "not-a-list-ok"}, TOK)  # tolerated
    frameworks.render_svg("nope", {}, TOK)  # unknown -> raises
except Exception:
    raised = True
check("unknown framework raises", raised)

# expand: writes SVG + replaces div; bad -> fallback, no crash; passthrough others.
with tempfile.TemporaryDirectory() as td:
    out = Path(td)
    doc = "Intro.\n\n::: funnel\nstages:\n  - {stage: A, value: 10}\n  - {stage: B, value: 4}\n:::\n\nOutro.\n"
    h = frameworks.expand(doc, "html", TOK, out)
    check("expand html img", "![](_framework-1.svg)" in h, h)
    check("expand wrote svg", (out / "_framework-1.svg").exists())
    check("expand prose kept", "Intro." in h and "Outro." in h)
    p = frameworks.expand(doc, "pdf", TOK, out)
    check("expand pdf image", "#image(" in p and "_framework-" in p, p)
    bad = frameworks.expand("::: matrix\n: : bad yaml :\n:::\n", "html", TOK, out)
    check("expand bad -> fallback", "could not render" in bad and "::: panel" in bad, bad[:80])
    other = "::: pullquote\nhi\n:::\n"
    check("expand passthrough", frameworks.expand(other, "html", TOK, out) == other)
    # pptx/other exports: no-op passthrough (PPTX is built natively by pptx_render, not here)
    check("expand non-linear passthrough", frameworks.expand(doc, "pptx", TOK, out) == doc)

# native PPTX shapes for all 6 types (pptx_render reuses frameworks parse helpers).
import tempfile as _tf  # noqa: E402

from studio import pptx_render  # noqa: E402

TOK_P = {"color": dict(TOK["color"], on_surface="#FFFFFF")}
PPTX_MD = (
    "# F\n\n::: funnel\nstages:\n  - {stage: A, value: 10}\n  - {stage: B, value: 4}\n:::\n\n"
    "## B\n\n::: bullseye\nrings:\n  - {ring: core, items: [X]}\n:::\n\n"
    "## M\n\n::: matrix\nitems:\n  - {label: I, x: low, y: high}\n:::\n\n"
    "## S\n\n::: swimlane\nlanes:\n  - {lane: L1, nodes: [A, B]}\n:::\n\n"
    "## D\n\n::: decision-tree\nroot: Q\nchildren:\n  - {condition: yes, root: P}\n:::\n\n"
    "## H\n\n::: heatmap\nrag: true\nrows: [r1]\ncols: [c1, c2]\ncells:\n  - [green, red]\n:::\n"
)
with _tf.TemporaryDirectory() as td2:
    out2 = Path(td2) / "d.pptx"
    pptx_render.build_pptx(PPTX_MD, TOK_P, out2)
    from pptx import Presentation  # noqa: E402

    sl = list(Presentation(str(out2)).slides)
    check("pptx 6 slides", len(sl) == 6, str(len(sl)))
    check("pptx slides have shapes", all(len(x.shapes) >= 2 for x in sl),
          str([len(x.shapes) for x in sl]))

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print("PASS: frameworks (bullseye/matrix/funnel/heatmap/swimlane/decision-tree)")
