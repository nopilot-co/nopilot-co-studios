#!/usr/bin/env python3
"""Parity matrix + fail-closed (ADR-006 / #129). Standalone; run: python3 tests/test_archetype_matrix.py

Driven by the capability registry (archetype_ir.CAPABILITIES): every declared
(archetype × backend) cell must render the archetype NATIVELY (no silent drop),
and an unknown fence must FAIL CLOSED — a visible placeholder + a stderr warning,
never a vanish. This is the regression gate for the render-convergence slice.
"""
from __future__ import annotations

import contextlib
import io
import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import archetype_ir as ir  # noqa: E402
from studio import gslide, pptx_render, uds_html  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


TOK = {"color": {"neutral": "#142F54", "surface": "#21456b", "tertiary": "#B14B33",
                 "on_primary": "#FFFFFF", "on_surface": "#FFFFFF", "secondary": "#6B7280", "primary": "#2A3548"}}

# A canonical sample body per archetype (each carries a "Discover"/category marker).
SAMPLES = {
    "chart": "type: bar\nx: [A, B, C]\ny: [3, 9, 6]",
    "flow": "- {title: Discover, caption: Audit}\n- {title: Design, caption: Wireframes}\n- {title: Build, caption: Ship}",
    "cards": "- {title: Discover, body: Audit}\n- {title: Design, body: Wireframes}\n- {title: Build, body: Ship}",
}


def _gslide_reqs(arch: str, body: str):
    md = f"---\ntitle: T\n---\n\n## H {{#h}}\n\n::: {arch}\n{body}\n:::\n"
    with tempfile.TemporaryDirectory() as td:
        src = Path(td) / "d.md"
        src.write_text(md)
        _t, reqs = gslide.build_requests(src, brand="nopilot")
    return reqs


def _pptx_shapes(arch: str, body: str):
    md = f"## H\n\n::: {arch}\n{body}\n:::\n"
    with tempfile.TemporaryDirectory() as td:
        out = Path(td) / "d.pptx"
        pptx_render.build_pptx(md, TOK, out)
        from pptx import Presentation

        return list(list(Presentation(str(out)).slides)[0].shapes)


def _html(arch: str, body: str) -> str:
    md = f"---\ntitle: T\n---\n\n## H {{#h}}\n\n::: {arch}\n{body}\n:::\n"
    _m, h = uds_html.render_body(md, brand="nopilot")
    return h


def _inserts(reqs) -> list[str]:
    return [r["insertText"]["text"] for r in reqs if "insertText" in r]


def _autoshapes(shapes) -> int:
    return len([x for x in shapes if "AUTO_SHAPE" in str(x.shape_type)])


# Per-archetype "rendered natively" predicate for each backend.
GSLIDE_OK = {
    "chart": lambda reqs: any("_bar" in r.get("updateShapeProperties", {}).get("objectId", "") for r in reqs),
    "flow": lambda reqs: "Discover" in _inserts(reqs),
    "cards": lambda reqs: "Discover" in _inserts(reqs),
}
PPTX_OK = {
    "chart": lambda sh: any(getattr(x, "has_chart", False) for x in sh),
    "flow": lambda sh: _autoshapes(sh) >= 3,
    "cards": lambda sh: _autoshapes(sh) >= 3,
}
HTML_OK = {
    "chart": lambda h: "uds-chart" in h and ":::" not in h,
    "flow": lambda h: "uds-flow" in h and ":::" not in h,
    "cards": lambda h: "uds-card__title" in h and ":::" not in h,
}
PROBE = {"gslide": (_gslide_reqs, GSLIDE_OK), "pptx": (_pptx_shapes, PPTX_OK), "html": (_html, HTML_OK)}

cells = 0
for arch, backends in sorted(ir.CAPABILITIES.items()):
    body = SAMPLES.get(arch)
    check(f"sample defined for '{arch}'", body is not None)
    if body is None:
        continue
    for bk in sorted(backends):
        render, ok = PROBE[bk]
        try:
            native = ok[arch](render(arch, body))
        except Exception as e:  # noqa: BLE001
            native = False
            check(f"matrix {arch} × {bk}: no crash", False, repr(e))
        check(f"matrix {arch} × {bk} renders native", native)
        cells += 1

# Fail-closed: an unknown fence is a visible placeholder + a stderr warning, not a drop.
md = "---\ntitle: T\n---\n\n## H {#h}\n\n::: zzz-unknown\nsome body here\n:::\n"
err = io.StringIO()
with contextlib.redirect_stderr(err):
    _m, html = uds_html.render_body(md, brand="nopilot")
check("fail-closed: visible placeholder", "unsupported block" in html, html[-200:])
check("fail-closed: stderr warned", "zzz-unknown" in err.getvalue(), err.getvalue())
check("fail-closed: no ::: leak", ":::" not in html)
check("fail-closed: content not silently dropped", "uds-muted" in html)

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print(f"PASS: parity matrix — {cells} (archetype×backend) cells native + fail-closed")
