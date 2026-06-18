#!/usr/bin/env python3
"""Cross-backend cards parity (ADR-006 / #129). Standalone; run: python3 tests/test_archetype_cards.py

ONE CardsNode reconciles the card dialects (bare list of {eyebrow,title,body} |
{cards:[…]} | strings); gslide / pptx / UDS-HTML each render every card natively.
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO / "design" / "scripts"))

from studio import archetype_ir as ir  # noqa: E402
from studio import gslide, pptx_render, uds_html  # noqa: E402

failures: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    if not cond:
        failures.append(f"{name}{(' — ' + detail) if detail else ''}")


# 1. dialects reconcile
lst = ir.normalise_cards([{"eyebrow": "01", "title": "Discover", "body": "Audit"}, {"title": "Design", "body": "Wireframes"}, {"title": "Build"}])
wrapped = ir.normalise_cards({"cards": [{"title": "Discover"}, {"title": "Design"}, {"title": "Build"}]})
check("list dialect → 3 cards", [c.title for c in lst.cards] == ["Discover", "Design", "Build"], str(lst.cards))
check("card eyebrow + body kept", (lst.cards[0].eyebrow, lst.cards[0].body) == ("01", "Audit"), str(lst.cards[0]))
check("{cards:[…]} dialect → 3", len(wrapped.cards) == 3)
check("bare strings → title-only", ir.normalise_cards(["A", "B"]).cards[0].title == "A")
check("malformed → CardsNode, no crash", isinstance(ir.normalise_cards("x: [bad"), ir.CardsNode))
check("capabilities[cards] = 3 backends", ir.CAPABILITIES["cards"] == {"gslide", "pptx", "html"})

CARDS_MD = ("---\ntitle: T\n---\n\n## Pillars {#p}\n\n::: cards\n"
            "- {eyebrow: Step one, title: Discover, body: Audit and interviews}\n"
            "- {eyebrow: Step two, title: Design, body: Wireframes and prototypes}\n"
            "- {eyebrow: Step three, title: Build, body: Implementation and QA}\n:::\n")
TOK = {"color": {"neutral": "#142F54", "surface": "#21456b", "tertiary": "#B14B33",
                 "on_primary": "#FFFFFF", "on_surface": "#FFFFFF", "secondary": "#6B7280", "primary": "#2A3548"}}

# 2. gslide native cards
with tempfile.TemporaryDirectory() as td:
    src = Path(td) / "d.md"
    src.write_text(CARDS_MD)
    _t, reqs = gslide.build_requests(src, brand="nopilot")
    inserts = [r["insertText"]["text"] for r in reqs if "insertText" in r]
    check("gslide: card titles present", all(t in inserts for t in ("Discover", "Design", "Build")), str(inserts))
    check("gslide: card bodies present", any("interviews" in t for t in inserts))

# 3. pptx native cards (new — pptx had no cards renderer before)
with tempfile.TemporaryDirectory() as td:
    out = Path(td) / "d.pptx"
    pptx_render.build_pptx("## Pillars\n\n::: cards\n- {title: Discover, body: x}\n- {title: Design, body: y}\n- {title: Build, body: z}\n:::\n", TOK, out)
    from pptx import Presentation

    sl = list(Presentation(str(out)).slides)[0]
    autoshapes = [sh for sh in sl.shapes if "AUTO_SHAPE" in str(sh.shape_type)]
    check("pptx: >= 3 card boxes", len(autoshapes) >= 3, str(len(autoshapes)))

# 4. UDS-HTML native cards (was a silent drop)
_m, html = uds_html.render_body(CARDS_MD, brand="nopilot")
check("html: uds-grid", 'class="uds-grid"' in html)
check("html: 3 native cards", html.count("uds-card__title") == 3, str(html.count("uds-card__title")))
check("html: titles present", all(t in html for t in ("Discover", "Design", "Build")))
check("html: bodies present", "interviews" in html and "Wireframes" in html)
check("html: fence consumed", ":::" not in html)

if failures:
    print(f"FAIL ({len(failures)})")
    for f in failures:
        print("  -", f)
    sys.exit(1)
print("PASS: cards archetype — native parity across gslide / pptx / UDS-HTML")
