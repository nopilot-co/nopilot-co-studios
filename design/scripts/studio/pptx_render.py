"""Editable PPTX engine (#19).

A parallel render path: builds a `.pptx` of NATIVE, editable PowerPoint shapes
(text boxes, autoshapes, native tables/charts) from the same `::: ` block source
the Quarto path uses — brand-tokenized from the resolved design tokens.

render.py forks here when the locked export is `pptx`. python-pptx is already a
dependency; this module imports it lazily so the module loads without it.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

# 16:9 deck geometry (EMU via Inches helper, applied lazily in build).
_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5

_DIV_RE = re.compile(
    r"^:::+\s*(?:\{\.)?(?P<name>[a-z][a-z0-9-]*)\}?\s*\n"
    r"(?P<body>.*?)\n"
    r"^:::+\s*$",
    re.MULTILINE | re.DOTALL,
)

# Blocks that begin a new slide on their own.
_SLIDE_BLOCKS = {"cover-slide", "section-slide"}


def _hex(color: str) -> str:
    return color.lstrip("#")


def split_slides(markdown: str) -> list[dict]:
    """Split source markdown into a list of slide dicts.

    A new slide starts at: a top-level `#`/`##` heading, or a `::: cover-slide` /
    `::: section-slide` block. Everything until the next such boundary is that
    slide's raw body. Each slide: {kind, title, body}. kind in
    {cover-slide, section-slide, content}.
    """
    lines = markdown.splitlines()
    slides: list[dict] = []
    cur: dict | None = None

    def _flush():
        nonlocal cur
        if cur is not None:
            cur["body"] = "\n".join(cur["_lines"]).strip()
            del cur["_lines"]
            slides.append(cur)
            cur = None

    i = 0
    while i < len(lines):
        line = lines[i]
        m_block = re.match(r"^:::+\s*(?:\{\.)?([a-z][a-z0-9-]*)\}?\s*$", line)
        if m_block and m_block.group(1) in _SLIDE_BLOCKS:
            _flush()
            name = m_block.group(1)
            body_lines = []
            i += 1
            while i < len(lines) and not re.match(r"^:::+\s*$", lines[i]):
                body_lines.append(lines[i])
                i += 1
            i += 1  # skip closing :::
            title = ""
            rest = []
            for bl in body_lines:
                h = re.match(r"^#+\s+(.*)$", bl)
                if h and not title:
                    title = h.group(1).strip()
                else:
                    rest.append(bl)
            slides.append(
                {"kind": name, "title": title, "body": "\n".join(rest).strip()}
            )
            continue
        h = re.match(r"^#{1,2}\s+(.*)$", line)
        if h:
            _flush()
            cur = {"kind": "content", "title": h.group(1).strip(), "_lines": []}
            i += 1
            continue
        if cur is None:
            cur = {"kind": "content", "title": "", "_lines": []}
        cur["_lines"].append(line)
        i += 1
    _flush()
    return [
        s
        for s in slides
        if s.get("title") or s.get("body") or s["kind"] in _SLIDE_BLOCKS
    ]


def build_pptx(markdown: str, tokens: dict[str, Any], out_path: Path) -> Path:
    """Build a .pptx from markdown + tokens, writing to out_path. Returns out_path."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    slides = split_slides(markdown)
    if not slides:
        # An empty source would yield a zero-slide deck, which some downstream
        # tools (e.g. the LibreOffice QA conversion) choke on. Emit one blank slide.
        slides = [{"kind": "content", "title": "", "body": ""}]
    for slide in slides:
        s = prs.slides.add_slide(blank)
        if slide["kind"] == "cover-slide":
            _render_cover(s, slide, tokens, prs)
        elif slide["kind"] == "section-slide":
            _render_section(s, slide, tokens, prs)
        else:
            _render_content(s, slide, tokens, prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _solid(shape, hex_color: str) -> None:
    from pptx.dml.color import RGBColor

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(_hex(hex_color))
    shape.line.fill.background()


def _set_text(tf, text: str, *, size: int, color: str, bold: bool = False) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(_hex(color))


def _render_cover(s, slide, tokens, prs) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(_SLIDE_W_IN), Inches(_SLIDE_H_IN)
    )
    _solid(bg, c["neutral"])
    title = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.4))
    _set_text(
        title.text_frame,
        slide.get("title") or "",
        size=44,
        color=c["on_primary"],
        bold=True,
    )
    if slide.get("body"):
        sub = s.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(1.0))
        _set_text(
            sub.text_frame,
            slide["body"].splitlines()[0],
            size=20,
            color=c["on_primary"],
        )


def _render_section(s, slide, tokens, prs) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(_SLIDE_W_IN), Inches(_SLIDE_H_IN)
    )
    _solid(bg, c["surface"])
    title = s.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.2))
    _set_text(
        title.text_frame,
        slide.get("title") or "",
        size=36,
        color=c["on_surface"],
        bold=True,
    )


def _render_content(s, slide, tokens, prs) -> None:
    from pptx.util import Inches

    c = tokens["color"]
    top = 0.4
    if slide.get("title"):
        t = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.9))
        _set_text(t.text_frame, slide["title"], size=30, color=c["primary"], bold=True)
        top = 1.5

    for block in _content_blocks(slide.get("body", "")):
        top = _place_block(s, block, tokens, top)


def _content_blocks(body: str) -> list[dict]:
    """Sequence a slide body into blocks: fenced divs, markdown tables, text runs."""
    blocks: list[dict] = []
    text_acc: list[str] = []

    def _flush_text():
        if text_acc:
            joined = "\n".join(text_acc).strip()
            if joined:
                blocks.append({"kind": "text", "body": joined})
            text_acc.clear()

    lines = body.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        m = re.match(r"^:::+\s*(?:\{\.)?([a-z][a-z0-9-]*)\}?\s*$", line)
        if m:
            _flush_text()
            name = m.group(1)
            inner = []
            i += 1
            while i < len(lines) and not re.match(r"^:::+\s*$", lines[i]):
                inner.append(lines[i])
                i += 1
            i += 1
            blocks.append({"kind": name, "body": "\n".join(inner).strip()})
            continue
        if line.strip().startswith("|") and "|" in line:
            _flush_text()
            tbl = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                tbl.append(lines[i])
                i += 1
            blocks.append({"kind": "table", "rows": _parse_table(tbl)})
            continue
        text_acc.append(line)
        i += 1
    _flush_text()
    return blocks


def _parse_table(lines: list[str]) -> list[list[str]]:
    rows = []
    for ln in lines:
        if re.match(r"^\s*\|?\s*:?-{2,}", ln):  # separator row
            continue
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        rows.append(cells)
    return rows


_PANEL_FILL = {
    "panel": None,
    "highlight": "surface",
    "ds-callout": "surface",
    "pullquote": None,
    "stat-panel": "surface",
    "kpi": "surface",
}


def _place_block(s, block, tokens, top: float) -> float:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    kind = block["kind"]
    if kind == "table":
        return _place_table(s, block["rows"], tokens, top)
    if kind == "kpi":
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6),
            Inches(top),
            Inches(6),
            Inches(1.6),
        )
        _solid(shp, c["surface"])
        _set_text(
            shp.text_frame, block["body"], size=40, color=c["tertiary"], bold=True
        )
        return top + 1.9
    if kind == "chart":
        return _place_chart(s, block["body"], tokens, top)
    if kind in ("cards", "card-grid"):
        from . import archetype_ir

        return _place_cards(s, archetype_ir.normalise_cards(block["body"]), tokens, top)
    if kind in ("flow", "process", "timeline", "hierarchy", "org"):
        return _place_diagram(s, kind, block["body"], tokens, top)
    if kind == "swimlane":
        return _place_swimlane(s, block["body"], tokens, top)
    if kind == "decision-tree":
        return _place_decision_tree(s, block["body"], tokens, top)
    if kind == "funnel":
        return _place_funnel(s, block["body"], tokens, top)
    if kind == "bullseye":
        return _place_bullseye(s, block["body"], tokens, top)
    if kind == "matrix":
        return _place_matrix(s, block["body"], tokens, top)
    if kind == "heatmap":
        return _place_heatmap(s, block["body"], tokens, top)
    if kind in _PANEL_FILL:
        h = 1.3
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(top), Inches(12), Inches(h)
        )
        fill = _PANEL_FILL[kind]
        if fill:
            _solid(shp, c[fill])
        else:
            shp.fill.background()
            from pptx.dml.color import RGBColor

            shp.line.color.rgb = RGBColor.from_string(_hex(c["secondary"]))
        # Surface-filled panels need on_surface for contrast; transparent panels
        # (panel/pullquote) sit on the slide background, so keep primary.
        text_color = c["on_surface"] if fill else c["primary"]
        _set_text(shp.text_frame, block["body"], size=18, color=text_color)
        return top + h + 0.3
    # default: text
    tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(1.2))
    _set_text_multiline(tb.text_frame, block["body"], tokens)
    return top + 1.3


def _place_table(s, rows, tokens, top: float) -> float:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    c = tokens["color"]
    if not rows:
        return top
    nrows, ncols = len(rows), max(len(r) for r in rows)
    gframe = s.shapes.add_table(
        nrows, ncols, Inches(0.6), Inches(top), Inches(8), Inches(0.4 * nrows)
    )
    tbl = gframe.table
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            cell.text = row[ci] if ci < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor.from_string(
                        # Header row sits on a surface fill → on_surface for contrast.
                        _hex(c["on_surface"] if ri == 0 else c["primary"])
                    )
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(_hex(c["surface"]))
    return top + 0.4 * nrows + 0.3


def _place_chart(s, body, tokens, top: float) -> float:
    from pptx.util import Inches

    try:
        from . import archetype_ir
        from . import charts as charts_mod

        node = archetype_ir.normalise_chart(body)  # one normaliser, shared across backends
        if node.chart_type not in charts_mod.CHART_TYPES:
            raise ValueError(f"unknown chart type '{node.chart_type}'")
        _add_native_chart(s, node, tokens, top)
        return top + 4.2
    except Exception:  # noqa: BLE001 — degrade: a chart must never crash the deck
        tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.8))
        _set_text(
            tb.text_frame,
            "[chart could not render]",
            size=14,
            color=tokens["color"]["secondary"],
        )
        return top + 1.0


_XL = {
    "bar": "COLUMN_CLUSTERED",
    "line": "LINE",
    "area": "AREA",
    "pie": "PIE",
    "scatter": "XY_SCATTER",
}


def _add_native_chart(s, node, tokens, top) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    c = tokens["color"]
    xltype = getattr(XL_CHART_TYPE, _XL.get(node.chart_type, "COLUMN_CLUSTERED"))
    cd = CategoryChartData()
    if node.chart_type == "pie":
        vals = node.series[0].values if node.series else []
        cd.categories = node.categories or [str(i) for i in range(len(vals))]
        cd.add_series("", tuple(vals))
    else:
        first = node.series[0].values if node.series else []
        cd.categories = node.categories or [str(i) for i in range(len(first))]
        for ser in node.series:
            cd.add_series(ser.name or "series", tuple(ser.values))
    gframe = s.shapes.add_chart(
        xltype, Inches(0.6), Inches(top), Inches(8.5), Inches(4.0), cd
    )
    chart = gframe.chart
    chart.has_legend = any(ser.name for ser in node.series)
    try:
        palette = [c["tertiary"], c["primary"], c["secondary"]]
        for plot in chart.plots:
            for si, ser in enumerate(plot.series):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = RGBColor.from_string(
                    _hex(palette[si % len(palette)])
                )
    except Exception:  # noqa: BLE001 — colour is best-effort; never fail the render
        pass


def _place_cards(s, node, tokens, top: float) -> float:
    """Native PPTX card grid from a CardsNode — rounded-rect cards (eyebrow? + title +
    body) laid out ≤3 across, wrapping to rows. Brand-tokenised; degrades to nothing
    on an empty node."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    c = tokens["color"]
    cards = node.cards
    if not cards:
        return top
    n = len(cards)
    cols = min(n, 3)
    gap, cw, ch = 0.3, (12.4 - 0.3 * (min(n, 3) - 1)) / min(n, 3), 1.6
    for i, card in enumerate(cards):
        x = 0.6 + (i % cols) * (cw + gap)
        y = top + (i // cols) * (ch + 0.3)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        _solid(shp, c["surface"])
        tf = shp.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        if card.eyebrow:
            r = para.add_run()
            r.text = card.eyebrow.upper()
            r.font.size, r.font.bold = Pt(9), True
            r.font.color.rgb = RGBColor.from_string(_hex(c["primary"]))
            para = tf.add_paragraph()
        rt = para.add_run()
        rt.text = card.title
        rt.font.size, rt.font.bold = Pt(15), True
        rt.font.color.rgb = RGBColor.from_string(_hex(c["on_surface"]))
        if card.body:
            pb = tf.add_paragraph()
            rb = pb.add_run()
            rb.text = card.body
            rb.font.size = Pt(11)
            rb.font.color.rgb = RGBColor.from_string(_hex(c["on_surface"]))
    rows = -(-n // cols)
    return top + rows * (ch + 0.3) + 0.2


def _place_diagram(s, kind, body, tokens, top: float) -> float:
    try:
        import yaml

        spec = yaml.safe_load(body) or {}
        if kind in ("flow", "process", "timeline"):
            if kind == "timeline":
                labels = [
                    str(e.get("label", ""))
                    for e in (spec.get("events") or [])
                    if isinstance(e, dict)
                ]
            else:
                from . import archetype_ir

                labels = [st.title for st in archetype_ir.normalise_flow(spec).steps]
            _diagram_linear(s, labels, tokens, top)
        else:
            from . import diagrams as diagrams_mod

            nodes, edges = diagrams_mod._flatten_tree(spec)
            _diagram_tree(s, nodes, edges, tokens, top)
        return top + 4.2
    except Exception:  # noqa: BLE001 — degrade
        from pptx.util import Inches

        tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.8))
        _set_text(
            tb.text_frame,
            f"[diagram '{kind}' could not render]",
            size=14,
            color=tokens["color"]["secondary"],
        )
        return top + 1.0


def _node_box(s, x_in, y_in, w_in, h_in, label, tokens):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in),
        Inches(y_in),
        Inches(w_in),
        Inches(h_in),
    )
    _solid(shp, c["neutral"])
    _set_text(shp.text_frame, label, size=14, color=c["on_primary"], bold=False)
    return shp


def _connect_v(s, a, b, tokens):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Emu

    c = tokens["color"]
    conn = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(a.left + a.width // 2),
        Emu(a.top + a.height),
        Emu(b.left + b.width // 2),
        Emu(b.top),
    )
    conn.line.color.rgb = RGBColor.from_string(_hex(c["tertiary"]))
    conn.line.width = Emu(19050)


def _connect_h(s, a, b, tokens):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Emu

    c = tokens["color"]
    conn = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(a.left + a.width),
        Emu(a.top + a.height // 2),
        Emu(b.left),
        Emu(b.top + b.height // 2),
    )
    conn.line.color.rgb = RGBColor.from_string(_hex(c["tertiary"]))
    conn.line.width = Emu(19050)


def _diagram_linear(s, labels, tokens, top: float) -> None:
    n = max(len(labels), 1)
    w = min(2.2, 11.0 / n - 0.3)
    gap = 0.3
    total = n * w + (n - 1) * gap
    x = max(0.6, (13.333 - total) / 2)
    boxes = []
    for lab in labels:
        boxes.append(_node_box(s, x, top + 1.4, w, 0.9, lab, tokens))
        x += w + gap
    for i in range(len(boxes) - 1):
        _connect_h(s, boxes[i], boxes[i + 1], tokens)


def _diagram_tree(s, nodes, edges, tokens, top: float) -> None:
    maxx = max((nd["x"] for nd in nodes), default=0) or 1
    maxd = max((nd["depth"] for nd in nodes), default=0) or 1
    w, h = 1.7, 0.7
    margin = 0.6
    # The x-coordinate positions the box's LEFT edge; the span must leave room for
    # a full box width on the right so the rightmost node stays on the 13.333" slide.
    span_w = max(_SLIDE_W_IN - 2 * margin - w, 1.0)
    row_h = 3.6 / (maxd + 1)
    placed = {}
    for nd in nodes:
        cx = margin + (nd["x"] / maxx) * span_w if maxx else margin
        cy = top + 1.0 + nd["depth"] * row_h
        placed[nd["id"]] = _node_box(s, cx, cy, w, h, nd["label"], tokens)
    for a, b in edges:
        _connect_v(s, placed[a], placed[b], tokens)


def _set_text_multiline(tf, text: str, tokens) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    c = tokens["color"]
    tf.word_wrap = True
    first = True
    for ln in text.splitlines():
        ln = ln.rstrip()
        if not ln:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = ("• " + ln[2:].strip()) if ln.strip().startswith("- ") else ln
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor.from_string(_hex(c["primary"]))


# --------------------------------------------------------------- framework shapes
# Native, editable PPTX shapes for the viz types whose authoring/parsing lives in
# studio.frameworks (reused here so the deck and the HTML/PDF SVG agree). Each
# degrades to a text note rather than crashing the deck (mirrors _place_chart).


def _degrade(s, label: str, tokens, top: float) -> float:
    from pptx.util import Inches

    tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.8))
    _set_text(tb.text_frame, f"[{label} could not render]", size=14,
              color=tokens["color"]["secondary"])
    return top + 1.0


def _line(s, x1, y1, x2, y2, tokens) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Emu, Inches

    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = RGBColor.from_string(_hex(tokens["color"]["secondary"]))
    conn.line.width = Emu(12700)


def _connect_centers(s, a, b, tokens) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Emu

    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Emu(a.left + a.width // 2), Emu(a.top + a.height // 2),
                                  Emu(b.left + b.width // 2), Emu(b.top + b.height // 2))
    conn.line.color.rgb = RGBColor.from_string(_hex(tokens["color"]["tertiary"]))
    conn.line.width = Emu(19050)


def _place_funnel(s, body, tokens, top: float) -> float:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    palette = [c["tertiary"], c["primary"], c["secondary"], c["neutral"], c["surface"]]
    try:
        import yaml

        spec = yaml.safe_load(body) or {}
        labels, vals = [], []
        for st in spec.get("stages") or []:
            if isinstance(st, dict):
                labels.append(str(st.get("stage", st.get("label", ""))))
                try:
                    vals.append(float(st.get("value", 0) or 0))
                except (TypeError, ValueError):
                    vals.append(0.0)
            else:
                labels.append(str(st))
                vals.append(0.0)
        n = len(labels) or 1
        maxv = max(vals) if any(vals) else 1.0
        row_h = min(0.7, 3.8 / n)
        full_w = 10.0
        y = top + 0.2
        for i, (lab, v) in enumerate(zip(labels, vals)):
            w = max((v / maxv) if maxv else 1.0, 0.12) * full_w
            shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6 + (full_w - w) / 2), Inches(y),
                                     Inches(w), Inches(row_h))
            _solid(shp, palette[i % len(palette)])
            vtxt = f"  ({int(v) if float(v).is_integer() else v})" if v else ""
            _set_text(shp.text_frame, f"{lab}{vtxt}", size=14, color=c["on_primary"])
            y += row_h + 0.12
        return y + 0.2
    except Exception:  # noqa: BLE001 — a viz must never crash the deck
        return _degrade(s, "funnel", tokens, top)


def _place_bullseye(s, body, tokens, top: float) -> float:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    c = tokens["color"]
    palette = [c["surface"], c["secondary"], c["neutral"], c["tertiary"], c["primary"]]
    try:
        import yaml

        from . import frameworks as fw

        spec = yaml.safe_load(body) or {}
        bands = fw._bands(spec) or [("", [])]
        n = len(bands)
        dia = 3.6
        cx, cy = 0.6 + dia / 2, top + 0.2 + dia / 2
        for j in range(n - 1, -1, -1):  # outermost first; core (bands[0]) on top
            d = dia * (j + 1) / n
            shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                     Inches(d), Inches(d))
            _solid(shp, palette[j % len(palette)])
        tx = 0.6 + dia + 0.4
        ty = top + 0.2
        for j, (label, items) in enumerate(bands):  # legend to the right (ovals overlap text)
            tb = s.shapes.add_textbox(Inches(tx), Inches(ty), Inches(_SLIDE_W_IN - tx - 0.6),
                                      Inches(0.6))
            _set_text(tb.text_frame, (label + ": " if label else "") + ", ".join(items),
                      size=14, color=c["primary"], bold=(j == 0))
            ty += 0.7
        return max(cy + dia / 2, ty) + 0.3
    except Exception:  # noqa: BLE001
        return _degrade(s, "bullseye", tokens, top)


def _place_matrix(s, body, tokens, top: float) -> float:
    from pptx.util import Inches

    c = tokens["color"]
    try:
        import yaml

        from . import frameworks as fw

        spec = yaml.safe_load(body) or {}
        axes = spec.get("axes") or {}
        ox, oy, sz = 1.2, top + 0.4, 4.0
        _line(s, ox + sz / 2, oy, ox + sz / 2, oy + sz, tokens)
        _line(s, ox, oy + sz / 2, ox + sz, oy + sz / 2, tokens)
        if axes.get("x"):
            tb = s.shapes.add_textbox(Inches(ox), Inches(oy + sz + 0.05), Inches(sz), Inches(0.3))
            _set_text(tb.text_frame, str(axes["x"]), size=12, color=c["secondary"])
        if axes.get("y"):
            tb = s.shapes.add_textbox(Inches(ox - 1.0), Inches(oy + sz / 2 - 0.15), Inches(1.0),
                                      Inches(0.3))
            _set_text(tb.text_frame, str(axes["y"]), size=12, color=c["secondary"])
        for it in fw._matrix_items(spec):
            px = ox + fw._pos(it.get("x")) * sz
            py = oy + (1 - fw._pos(it.get("y"))) * sz  # invert: high y at top
            tb = s.shapes.add_textbox(Inches(px - 0.7), Inches(py - 0.15), Inches(1.4), Inches(0.3))
            _set_text(tb.text_frame, str(it.get("label", it.get("name", ""))), size=11,
                      color=c["tertiary"], bold=True)
        return oy + sz + 0.5
    except Exception:  # noqa: BLE001
        return _degrade(s, "matrix", tokens, top)


def _place_swimlane(s, body, tokens, top: float) -> float:
    from pptx.util import Inches

    c = tokens["color"]
    try:
        import yaml

        from . import frameworks as fw

        spec = yaml.safe_load(body) or {}
        lanes = fw._lanes(spec) or [("", [])]
        nlanes = len(lanes)
        maxcols = max((len(nodes) for _, nodes in lanes), default=1) or 1
        lane_h = min(1.1, 4.2 / nlanes)
        label_w = 1.5
        area_x = 0.6 + label_w
        cell_w = min(2.0, (_SLIDE_W_IN - area_x - 0.6) / maxcols - 0.2) if maxcols else 2.0
        boxes: dict[str, Any] = {}
        flat: list[str] = []
        y0 = top + 0.3
        for li, (name, nodes) in enumerate(lanes):
            y = y0 + li * lane_h
            lbl = s.shapes.add_textbox(Inches(0.6), Inches(y + lane_h / 2 - 0.3),
                                       Inches(label_w - 0.1), Inches(0.6))
            _set_text(lbl.text_frame, name, size=12, color=c["primary"], bold=True)
            if li > 0:
                _line(s, 0.6, y, _SLIDE_W_IN - 0.6, y, tokens)
            for ci, txt in enumerate(nodes):
                box = _node_box(s, area_x + ci * (cell_w + 0.2), y + lane_h / 2 - 0.3,
                                cell_w, 0.6, txt, tokens)
                boxes[txt] = box
                flat.append(txt)
        edges = spec.get("edges")
        if isinstance(edges, list) and edges:
            pairs = [(str(e.get("from", "")), str(e.get("to", ""))) for e in edges
                     if isinstance(e, dict)]
        else:
            pairs = [(flat[i], flat[i + 1]) for i in range(len(flat) - 1)]
        for a, b in pairs:
            if a in boxes and b in boxes:
                _connect_centers(s, boxes[a], boxes[b], tokens)
        return y0 + nlanes * lane_h + 0.3
    except Exception:  # noqa: BLE001
        return _degrade(s, "swimlane", tokens, top)


def _place_decision_tree(s, body, tokens, top: float) -> float:
    from pptx.util import Inches

    c = tokens["color"]
    try:
        import yaml

        from . import frameworks as fw

        spec = yaml.safe_load(body) or {}
        nodes, edges = fw._walk_decision(spec)
        maxx = max((nd["x"] for nd in nodes), default=0) or 1
        maxd = max((nd["depth"] for nd in nodes), default=0) or 1
        w, h, margin = 1.8, 0.6, 0.6
        span_w = max(_SLIDE_W_IN - 2 * margin - w, 1.0)
        row_h = min(1.3, 4.0 / (maxd + 1))
        placed = {}
        for nd in nodes:
            cx = margin + (nd["x"] / maxx) * span_w if maxx else margin
            cy = top + 0.4 + nd["depth"] * row_h
            box = _node_box(s, cx, cy, w, h, nd["label"], tokens)
            if nd["kind"] == "outcome":
                _solid(box, c["tertiary"])  # text already set by _node_box (on_primary)
            placed[nd["id"]] = box
        for a, b, cond in edges:
            _connect_v(s, placed[a], placed[b], tokens)
            if cond:
                ax, bx = placed[a], placed[b]
                midx = (ax.left + ax.width // 2 + bx.left + bx.width // 2) // 2
                midy = (ax.top + ax.height + bx.top) // 2
                tb = s.shapes.add_textbox(midx - Inches(0.6), midy - Inches(0.13),
                                          Inches(1.2), Inches(0.26))
                _set_text(tb.text_frame, cond, size=9, color=c["primary"])
        return top + 0.4 + (maxd + 1) * row_h + 0.3
    except Exception:  # noqa: BLE001
        return _degrade(s, "decision-tree", tokens, top)


def _place_heatmap(s, body, tokens, top: float) -> float:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    c = tokens["color"]
    try:
        import yaml

        from . import frameworks as fw

        spec = yaml.safe_load(body) or {}
        rag = bool(spec.get("rag"))
        rows, cols, grid = fw._heatmap_grid(spec)
        nr, nc = len(rows), len(cols)
        trows, tcols = nr + 1, nc + 1
        th = min(0.5 * trows, 4.6)
        gframe = s.shapes.add_table(trows, tcols, Inches(0.6), Inches(top + 0.2),
                                    Inches(min(2 + 1.6 * nc, 11)), Inches(th))
        tbl = gframe.table
        tbl.cell(0, 0).text = ""
        for ci, co in enumerate(cols):
            tbl.cell(0, ci + 1).text = str(co)
        for ri, rname in enumerate(rows):
            tbl.cell(ri + 1, 0).text = str(rname)
            for ci in range(nc):
                cell = tbl.cell(ri + 1, ci + 1)
                v = grid[ri][ci]
                cell.text = str(v)
                fc = fw._RAG.get(str(v).strip().lower()) if rag else None
                if fc:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(_hex(fc))
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(12)
                        run.font.color.rgb = RGBColor.from_string(
                            _hex(c["on_primary"] if fc else c["primary"]))
        for ci in range(tcols):  # header band
            hc = tbl.cell(0, ci)
            hc.fill.solid()
            hc.fill.fore_color.rgb = RGBColor.from_string(_hex(c["surface"]))
            for para in hc.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string(_hex(c["on_surface"]))
        return top + 0.2 + th + 0.3
    except Exception:  # noqa: BLE001
        return _degrade(s, "heatmap", tokens, top)
