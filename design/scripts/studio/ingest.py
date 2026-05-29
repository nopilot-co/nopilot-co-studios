"""Brand ingestion: deterministic extraction from source materials.

This module does NOT synthesize the brand spec — that's the LLM's job in the
brand-ingest skill. We extract raw facts (colors, fonts, master slides) and
write an _ingest-report.md plus a draft _brand.yml that the LLM then refines.

Source paths may be files OR folders. Folders are recursively walked and every
supported file inside is extracted; the folder itself is also copied as-is to
ingest-sources/ so the original structure is preserved for later re-ingestion.
"""

from __future__ import annotations

import shutil
from collections import Counter
from pathlib import Path
from typing import Iterable

import yaml

from . import TEMPLATES
from . import brand as brand_mod

# Extensions we know how to extract from. Anything else is copied but not analysed.
_SUPPORTED_EXTS = {".pdf", ".pptx", ".png", ".jpg", ".jpeg", ".svg"}

# Names/patterns to skip when walking folders. Hidden files (.DS_Store, .git/, etc.)
# are skipped by the leading-dot check.
_SKIP_DIRS = {"__pycache__", "node_modules", ".git", ".venv"}


def _walk_supported(path: Path) -> Iterable[Path]:
    """Yield files we know how to extract from.

    If `path` is a file, yields it (regardless of extension — caller filters).
    If `path` is a directory, walks recursively, yielding every file with a
    supported extension. Skips hidden files/dirs and a few well-known noise
    directories. Does not follow symlinks.
    """
    if path.is_file():
        yield path
        return
    for p in sorted(path.rglob("*")):
        if p.is_symlink():
            continue
        # Skip if any path component is hidden or a noise dir
        if any(
            part.startswith(".") or part in _SKIP_DIRS
            for part in p.relative_to(path).parts
        ):
            continue
        if p.is_file() and p.suffix.lower() in _SUPPORTED_EXTS:
            yield p


def run(slug: str, sources: list[Path]) -> str:
    """Run extraction on sources. Returns a markdown report."""
    # Writes to the resolved brand root — the shared studios-level store for new
    # brands, or the legacy location if this brand predates elevation.
    brand = brand_mod.brand_root(slug)
    (brand / "ingest-sources").mkdir(parents=True, exist_ok=True)
    (brand / "assets").mkdir(parents=True, exist_ok=True)
    (brand / "css").mkdir(parents=True, exist_ok=True)

    report_lines: list[str] = [f"# Ingest report — {slug}", ""]
    extracted_colors: list[str] = []
    extracted_fonts: set[str] = set()
    has_pptx = False
    logo_paths: list[Path] = []

    for src in sources:
        # 1. Preserve the original under ingest-sources/, file or whole folder
        dest = brand / "ingest-sources" / src.name
        if src.is_file():
            shutil.copy2(src, dest)
        else:
            shutil.copytree(src, dest, dirs_exist_ok=True, symlinks=False)

        # 2. Walk the source for supported files (a folder may yield 0..N)
        files = list(_walk_supported(src))
        report_lines.append(f"## {src.name}" + ("/" if src.is_dir() else ""))
        if src.is_dir():
            if not files:
                report_lines.append("- (no supported files found inside)")
                report_lines.append("")
                continue
            report_lines.append(
                f"- Walked folder, found {len(files)} supported file(s):"
            )

        for f in files:
            rel = f.relative_to(src) if src.is_dir() else f.name
            ext = f.suffix.lower()
            label = f"  - `{rel}`" if src.is_dir() else f"- `{f.name}`"
            if ext == ".pdf":
                colors, fonts = _extract_from_pdf(f)
                extracted_colors.extend(colors)
                extracted_fonts.update(fonts)
                report_lines.append(
                    f"{label} (pdf) — colors: {', '.join(colors[:5]) or '(none)'} · "
                    f"fonts: {', '.join(sorted(fonts)) or '(none)'}"
                )
            elif ext == ".pptx":
                has_pptx = True
                colors, fonts, ref_path = _extract_from_pptx(f, brand)
                extracted_colors.extend(colors)
                extracted_fonts.update(fonts)
                report_lines.append(
                    f"{label} (pptx) — colors: {', '.join(colors[:5]) or '(none)'} · "
                    f"fonts: {', '.join(sorted(fonts)) or '(none)'} · "
                    f"reference deck → {ref_path.name}"
                )
            elif ext in (".png", ".jpg", ".jpeg", ".svg"):
                logo_paths.append(f)
                colors = _extract_from_image(f)
                extracted_colors.extend(colors)
                report_lines.append(
                    f"{label} (image) — dominant colors: "
                    f"{', '.join(colors[:5]) or '(none)'}"
                )
        report_lines.append("")

    # Promote first logo to canonical
    if logo_paths:
        primary_logo = logo_paths[0]
        ext = primary_logo.suffix.lower()
        dest_logo = brand / "assets" / f"logo{ext}"
        shutil.copy2(primary_logo, dest_logo)
        report_lines.append(f"**Logo:** copied {primary_logo.name} → assets/logo{ext}")
        report_lines.append("")

    # Tally colors
    color_counts = Counter(c.upper() for c in extracted_colors if c)
    top_colors = [c for c, _ in color_counts.most_common(8)]

    # Write a draft _brand.yml
    draft = _draft_brand_yml(slug, top_colors, sorted(extracted_fonts))
    brand_yml = brand / "_brand.yml"
    if brand_yml.exists():
        backup = brand / "_brand.yml.bak"
        shutil.copy2(brand_yml, backup)
        report_lines.append(f"**Existing _brand.yml backed up to:** {backup}")
    brand_yml.write_text(
        yaml.safe_dump(draft, sort_keys=False, default_flow_style=False)
    )

    # Stub voice/style files if absent
    for fname, content in [
        ("tone-of-voice.md", _voice_stub(slug)),
        ("style-guide.md", _style_stub(slug)),
    ]:
        p = brand / fname
        if not p.exists():
            p.write_text(content)

    report_lines.extend(
        [
            "## Draft _brand.yml written",
            "",
            f"Path: `{brand_yml}`",
            "",
            "**Decisions needed from the LLM (brand-ingest skill):**",
            "- Confirm/correct primary color (currently set to most-frequent extracted color)",
            "- Choose heading vs body typeface from the extracted font list",
            "- Add accent/dataviz colors if needed",
            "- Fill `tone-of-voice.md` and `style-guide.md` from source guideline text",
        ]
    )
    if not has_pptx:
        report_lines.append(
            "- No PPTX source — run `studio ingest synthesize-pptx --brand "
            f"{slug}` to generate a reference deck from the brand tokens"
        )

    report = "\n".join(report_lines)
    (brand / "_ingest-report.md").write_text(report)
    return report


# -------------------------------------------------------- extractors


def _extract_from_pdf(pdf: Path) -> tuple[list[str], set[str]]:
    """Extract dominant colors and font family names from a PDF."""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return [], set()

    fonts: set[str] = set()
    color_samples: list[str] = []

    doc = pdfium.PdfDocument(str(pdf))
    # Sample only first 5 pages — guidelines PDFs are typically front-loaded
    for i in range(min(5, len(doc))):
        page = doc[i]
        bitmap = page.render(scale=1.0)
        pil = bitmap.to_pil().convert("RGB")
        color_samples.extend(_dominant_colors(pil, n=5))
    return color_samples, fonts  # font extraction from PDF is unreliable; leave to LLM


def _extract_from_pptx(pptx: Path, brand_dir: Path) -> tuple[list[str], set[str], Path]:
    """Extract colors + fonts from PPTX and copy as reference deck."""
    from pptx import Presentation

    ref_dest = brand_dir / "reference.pptx"
    shutil.copy2(pptx, ref_dest)

    prs = Presentation(str(pptx))
    fonts: set[str] = set()
    colors: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        try:
                            rgb = run.font.color.rgb
                            if rgb:
                                colors.append(f"#{str(rgb)}")
                        except (AttributeError, TypeError):
                            pass
    return colors, fonts, ref_dest


def _extract_from_image(img_path: Path) -> list[str]:
    """Dominant colors from a logo image."""
    try:
        from PIL import Image
    except ImportError:
        return []
    pil = Image.open(img_path).convert("RGBA")
    # drop fully-transparent pixels
    pixels = [(r, g, b) for r, g, b, a in pil.getdata() if a > 200]
    if not pixels:
        return []
    from PIL import Image as _Image

    tmp = _Image.new("RGB", (len(pixels), 1))
    tmp.putdata(pixels)
    return _dominant_colors(tmp, n=5)


def _dominant_colors(pil_image, n: int = 5) -> list[str]:
    """k-means-ish via PIL's adaptive palette."""
    from PIL import Image

    img = pil_image.convert("RGB").resize((128, 128))
    pal = img.convert("P", palette=Image.Palette.ADAPTIVE, colors=n)
    palette = pal.getpalette()
    counts = sorted(pal.getcolors(), reverse=True)
    hexes: list[str] = []
    for _, idx in counts[:n]:
        r, g, b = palette[idx * 3 : idx * 3 + 3]
        # Skip near-white / near-black (usually paper/ink, not brand)
        if max(r, g, b) > 245 and min(r, g, b) > 245:
            continue
        if max(r, g, b) < 15:
            continue
        hexes.append(f"#{r:02X}{g:02X}{b:02X}")
    return hexes


# -------------------------------------------------------- drafts


def _draft_brand_yml(slug: str, colors: list[str], fonts: list[str]) -> dict:
    primary = colors[0] if colors else "#0066CC"
    secondary = colors[1] if len(colors) > 1 else "#333333"
    headings_font = fonts[0] if fonts else "Inter"
    body_font = fonts[1] if len(fonts) > 1 else headings_font
    return {
        "meta": {"name": slug, "link": ""},
        "color": {
            "palette": {f"brand-{i+1}": c for i, c in enumerate(colors[:6])},
            "primary": primary,
            "secondary": secondary,
            "foreground": "#1A1A1A",
            "background": "#FFFFFF",
        },
        "typography": {
            "fonts": [
                {"family": headings_font, "source": "google"},
                {"family": body_font, "source": "google"},
            ],
            "base": {"family": body_font, "size": "11pt"},
            "headings": {"family": headings_font, "weight": 600},
            "monospace": {"family": "JetBrains Mono", "source": "google"},
        },
        "logo": {
            "small": "assets/logo.svg",
            "medium": "assets/logo.svg",
            "large": "assets/logo.svg",
        },
    }


def _voice_stub(slug: str) -> str:
    return f"""# Tone of voice — {slug}

> This file is injected into Claude's context during composition. It is not rendered into outputs.

## Attributes
- (to be written from source guidelines)

## Forbidden
- (e.g. "synergy", "leverage as a verb", "we're excited to announce")

## Preferred
- (e.g. "use 'helps you' rather than 'enables'")

## Example
> (paste a 2–3 sentence passage that exemplifies the voice)
"""


def _style_stub(slug: str) -> str:
    return f"""# Style guide — {slug}

> Writing mechanics. Not rendered.

## Capitalisation
- Headings: sentence case
- Brand name: as written

## Numbers & dates
- Spell out one to ten in prose; use numerals from 11
- Dates: 25 May 2026

## Punctuation
- Oxford comma: yes / no

## Glossary
- (preferred terms and what they replace)
"""


def synthesize_reference_pptx(slug: str) -> Path:
    """Render an empty 5-slide deck via Quarto to use as reference.pptx."""
    import tempfile

    brand = brand_mod.brand_root(slug)
    if not (brand / "_brand.yml").exists():
        raise FileNotFoundError(f"no _brand.yml for {slug}")

    sample_md = TEMPLATES / "pptx" / "reference-skeleton.md"
    if not sample_md.exists():
        # Fallback: inline a minimal 5-slide skeleton
        sample_md_content = """---
title: "Reference deck"
---

# Title slide

::: notes
Cover slide.
:::

## Section divider

## Content slide

- Point one
- Point two
- Point three

## Two-column layout

::: {.columns}
::: {.column width="50%"}
Left
:::
::: {.column width="50%"}
Right
:::
:::

## Closing
"""
        with tempfile.TemporaryDirectory() as td:
            td_path = Path(td)
            (td_path / "sample.md").write_text(sample_md_content)
            shutil.copy2(brand / "_brand.yml", td_path / "_brand.yml")
            (td_path / "_quarto.yml").write_text(
                "project:\n  type: default\nformat:\n  pptx: default\nbrand: _brand.yml\n"
            )
            import subprocess

            result = subprocess.run(
                ["quarto", "render", "sample.md", "--to", "pptx"],
                cwd=td_path,
                capture_output=True,
                text=True,
            )
            if result.returncode != 0:
                raise RuntimeError(f"synthesize failed:\n{result.stderr}")
            produced = td_path / "sample.pptx"
            dest = brand / "reference.pptx"
            shutil.move(str(produced), dest)
            return dest

    raise NotImplementedError("custom sample.md path not yet wired")
